import pandas as pd
from flask import Flask, request, render_template, send_file
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches

app = Flask(__name__)

import matplotlib.pyplot as plt

def create_chart(data):
    # Esempio di creazione di un grafico utilizzando pandas e matplotlib
    # Puoi personalizzare questa funzione in base alle tue esigenze
    # data: DataFrame di pandas contenente i dati da plottare
    # return: Oggetto figura matplotlib

    # Creazione di una nuova figura specificando il backend "Agg"
    plt.figure(figsize=(8, 6), tight_layout=True, dpi=100)

    # Esempio di creazione di un grafico a barre
    data.plot(kind='bar')
    plt.title('Grafico')
    plt.xlabel('X-axis')
    plt.ylabel('Y-axis')

    # Restituisci l'oggetto figura
    return plt.gcf()

def create_ppt(chart, start_date, end_date):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Layout del titolo e contenuto

    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = f"Grafico dal {start_date} al {end_date}"

    img_buffer = BytesIO()
    chart.savefig(img_buffer)
    img_buffer.seek(0)
    left = Inches(2)
    top = Inches(2)
    pic = slide.shapes.add_picture(img_buffer, left, top)

    return prs

@app.route('/', methods=['GET'])
def index():
    return render_template('index.html')

@app.route('/process_excel', methods=['POST'])
def process_excel():
    excel_file = request.files['excelFile']
    start_date = request.form['startDate']
    end_date = request.form['endDate']

    # Carica il file Excel e crea un DataFrame con pandas
    df = pd.read_excel(excel_file)

    # Converti le colonne contenenti dati temporali in formato datetime di pandas
    df['DATA'] = pd.to_datetime(df['DATA'])

    # Effettua l'elaborazione dei dati e crea il grafico
    chart = create_chart(df)

    # Crea il PowerPoint e aggiunge il grafico
    prs = create_ppt(chart, start_date, end_date)

    # Salva il PowerPoint in un buffer di byte
    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)

    return send_file(ppt_buffer, as_attachment=True, attachment_filename='grafici_powerpoint.pptx')

if __name__ == '__main__':
    app.run(debug=True)
